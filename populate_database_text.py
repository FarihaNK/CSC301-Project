"""
populate_database_text.py

Populates a Chroma DB with *text* documents only (PDF, TXT, DOCX, PPTX),
using your MPNet embedding model.
"""
import argparse
import os
import time
import shutil
from langchain_community.document_loaders import PyPDFLoader
from langchain.schema.document import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from get_embedding import get_embedding_function
from langchain.text_splitter import SpacyTextSplitter
from docx import Document as DocxDocument
from langchain_chroma import Chroma
from pptx import Presentation

CHROMA_PATH = "./chroma_db"

def main(data_path):
    parser = argparse.ArgumentParser()
    parser.add_argument("--reset", action="store_true", help="Reset the database.")
    args = parser.parse_args()

    if args.reset:
        print("Clearing Database...")
        clear_database()

    # Load *only* text-based documents
    print("Loading text documents (pdf/txt/docx/pptx)")
    documents = load_text_documents(data_path)
    print(f"Loaded {len(documents)} total text documents")

    # Chunk text documents
    chunks = split_documents(documents)
    print(f"Split into {len(chunks)} text chunks")

    # Add them to the text-based DB
    add_to_chroma(chunks)
    print(f"Added {len(chunks)} text chunks into Chroma DB successfully.")

def load_text_documents(data_path):
    """Loads ONLY PDF, TXT, DOCX, PPTX. No images."""
    all_docs = []

    pdf_docs = load_pdfs_recursive(data_path)
    all_docs.extend(pdf_docs)

    txt_docs = load_txt_files_recursive(data_path)
    all_docs.extend(txt_docs)

    docx_docs = load_docx_files_recursive(data_path)
    all_docs.extend(docx_docs)

    pptx_docs = load_pptx_files_recursive(data_path)
    all_docs.extend(pptx_docs)

    return all_docs

def load_pdfs_recursive(directory):
    docs = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".pdf"):
                file_path = os.path.join(root, file)
                loader = PyPDFLoader(file_path)
                try:
                    new_docs = loader.load()
                    docs.extend(new_docs)
                except Exception as e:
                    print(f"Error loading PDF file {file_path}: {e}")
    return docs

def load_txt_files_recursive(directory):
    docs = []
    encodings = ['iso-8859-1','latin-1','utf-8']
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".txt"):
                file_path = os.path.join(root, file)
                content = ""
                for enc in encodings:
                    try:
                        with open(file_path,'r',encoding=enc) as f:
                            content = f.read()
                        break
                    except Exception as e:
                        pass
                if content:
                    docs.append(
                        Document(
                            page_content=content,
                            metadata={"source": file_path,"title":os.path.splitext(file)[0]}
                        )
                    )
    return docs

def load_docx_files_recursive(directory):
    docs = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".docx"):
                file_path = os.path.join(root, file)
                try:
                    doc = DocxDocument(file_path)
                    text = "\n".join([p.text for p in doc.paragraphs])
                    docs.append(
                        Document(
                            page_content=text,
                            metadata={"source": file_path,"title":os.path.splitext(file)[0]}
                        )
                    )
                except Exception as e:
                    print(f"Error loading DOCX file {file_path}: {e}")
    return docs

def load_pptx_files_recursive(directory):
    docs = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".pptx"):
                file_path = os.path.join(root, file)
                try:
                    prs = Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text.append(shape.text)
                    combined = "\n".join(text)
                    docs.append(
                        Document(
                            page_content=combined,
                            metadata={"source": file_path,"title":os.path.splitext(file)[0]}
                        )
                    )
                except Exception as e:
                    print(f"Error loading PPTX file {file_path}: {e}")
    return docs

def split_documents(documents):
    """
    Split documents by sentences using SpaCy,
    respecting semantic boundaries rather than raw character counts.
    """
    text_splitter = SpacyTextSplitter(
        pipeline='en_core_web_sm',
        chunk_size=1000,
        chunk_overlap=50
    )
    all_chunks = []
    for doc in documents:
        text_splits = text_splitter.split_text(doc.page_content)
        for chunk_str in text_splits:
            chunk_doc = Document(
                page_content=chunk_str,
                metadata=doc.metadata
            )
            all_chunks.append(chunk_doc)
    return all_chunks

def get_database():
    return Chroma(
        persist_directory=CHROMA_PATH,
        embedding_function=get_embedding_function()  # MPNet-based
    )

def add_to_chroma(chunks):
    db = get_database()
    db.add_documents(chunks)

def clear_database():
    if os.path.exists(CHROMA_PATH):
        shutil.rmtree(CHROMA_PATH)

if __name__ == "__main__":
    data_path = "./sampledata"
    main(data_path)
