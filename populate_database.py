# import argparse
# import os
# import time
# import shutil
# import warnings
# from langchain_community.document_loaders import PyPDFDirectoryLoader
# from langchain_community.document_loaders import PyPDFLoader
# from langchain.schema.document import Document
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# from get_embedding import get_embedding_function
# from docx import Document as DocxDocument
# #from langchain.vectorstores import Chroma
# from langchain_chroma import Chroma
# from pptx import Presentation
# import pandas as pd

# # Suppress specific warnings
# # warnings.filterwarnings("ignore", message="Cannot parse header or footer so it will be ignored")
# # warnings.filterwarnings("ignore", message="Data Validation extension is not supported and will be removed")
# # warnings.filterwarnings("ignore", message="Conditional Formatting extension is not supported and will be removed")
# # warnings.filterwarnings("ignore", message="Unknown extension is not supported and will be removed")
# # warnings.filterwarnings("ignore", message="Slicer List extension is not supported and will be removed")

# CHROMA_PATH = "./chroma_db"

# #db = Chroma(persist_directory=CHROMA_PATH, embedding_function=get_embedding_function())


# def main(data_path):
#     parser = argparse.ArgumentParser()
#     parser.add_argument("--reset", action="store_true", help="Reset the database.")
#     args = parser.parse_args()

#     if args.reset:
#         print("Clearing Database")
#         clear_database()

#     print("Loading documents")
#     documents = load_documents(data_path)
#     print(f"Loaded {len(documents)} documents")
#     chunks = split_documents(documents)
#     print(f"Split into {len(chunks)} chunks")
#     add_to_chroma(chunks)
#     print(f"added {len(chunks)} into Chroma DB Sucessfully :)")

# def load_documents(data_path):
#     all_documents = []

#     print("Loading PDF documents")
#     pdf_documents = load_pdfs_recursive(data_path)
#     all_documents.extend(pdf_documents)
#     print(f"Loaded {len(pdf_documents)} PDF documents")

#     print("Loading TXT documents")
#     txt_documents = load_txt_files_recursive(data_path)
#     all_documents.extend(txt_documents)
#     print(f"Loaded {len(txt_documents)} TXT documents")

#     print("Loading DOCX documents")
#     docx_documents = load_docx_files_recursive(data_path)
#     all_documents.extend(docx_documents)
#     print(f"Loaded {len(docx_documents)} DOCX documents")

#     print("Loading PPTX documents")
#     pptx_documents = load_pptx_files_recursive(data_path)
#     all_documents.extend(pptx_documents)
#     print(f"Loaded {len(pptx_documents)} PPTX documents")

#     return all_documents

# def load_pdfs_recursive(directory):
#     pdf_documents = []
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.lower().endswith(".pdf"):
#                 file_path = os.path.join(root, file)
#                 loader = PyPDFLoader(file_path)
#                 try:
#                     docs = loader.load()
#                     pdf_documents.extend(docs)
#                 except Exception as e:
#                     print(f"Error loading PDF file {file_path}: {e}")
#     return pdf_documents

# def load_txt_files_recursive(directory):
#     txt_documents = []
#     encodings = ['iso-8859-1', 'latin-1']
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.endswith(".txt"):
#                 file_path = os.path.join(root, file)
#                 content = ""
#                 for encoding in encodings:
#                     try:
#                         with open(file_path, 'r', encoding=encoding) as f:
#                             content = f.read()
#                         break
#                     except (UnicodeDecodeError, IOError) as e:
#                         print(f"Failed to read {file_path} with encoding {encoding}: {e}")
#                 txt_documents.append(Document(page_content=content, metadata={"source": file_path, "title": os.path.splitext(file)[0]}))
#     return txt_documents

# def load_docx_files_recursive(directory):
#     docx_documents = []
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.endswith(".docx"):
#                 file_path = os.path.join(root, file)
#                 try:
#                     doc = DocxDocument(file_path)
#                     content = "\n".join([para.text for para in doc.paragraphs])
#                     docx_documents.append(Document(page_content=content, metadata={"source": file_path, "title": os.path.splitext(file)[0]}))
#                 except Exception as e:
#                     print(f"Error loading DOCX file {file_path}: {e}")
#     return docx_documents

# def load_pptx_files_recursive(directory):
#     pptx_documents = []
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.endswith(".pptx"):
#                 file_path = os.path.join(root, file)
#                 try:
#                     prs = Presentation(file_path)
#                     text = []
#                     for slide in prs.slides:
#                         for shape in slide.shapes:
#                             if shape.has_text_frame:
#                                 text.append(shape.text)
#                     pptx_documents.append(Document(page_content='\n'.join(text), metadata={"source": file_path, "title": os.path.splitext(file)[0]}))
#                 except Exception as e:
#                     print(f"Error loading PPTX file {file_path}: {e}")
#     return pptx_documents

# def split_documents(documents):
#     text_splitter = RecursiveCharacterTextSplitter(
#         chunk_size=1000,
#         chunk_overlap=10,
#         length_function=len,
#         is_separator_regex=False
#     )
#     chunks = text_splitter.split_documents(documents)
#     return chunks


# def get_database():
#     return Chroma(persist_directory=CHROMA_PATH, embedding_function=get_embedding_function())

# def add_to_chroma(chunks):
    
#     chunks_with_ids = calculate_chunk_ids(chunks)
#     db = get_database()
#     db.add_documents(chunks_with_ids)

# def calculate_chunk_ids(chunks):
#     for chunk in chunks:
#         chunk.metadata["id"] = f'{chunk.metadata.get("source")}'
#     return chunks

# def clear_database():
#     if os.path.exists(CHROMA_PATH):
#         shutil.rmtree(CHROMA_PATH)
#         print("Database cleared")

# if __name__ == "__main__":
#     DATA_PATH="./sampledata"
#     main(DATA_PATH)



#STABLE VERSION FOR TEXT BASED.

# import argparse
# import os
# import time
# import shutil
# import warnings
# from langchain_community.document_loaders import PyPDFDirectoryLoader
# from langchain_community.document_loaders import PyPDFLoader
# from langchain.schema.document import Document
# from langchain_text_splitters import RecursiveCharacterTextSplitter
# from get_embedding import get_embedding_function
# from langchain.text_splitter import SpacyTextSplitter

# from docx import Document as DocxDocument
# from langchain_chroma import Chroma
# from pptx import Presentation
# import pandas as pd

# # Optional: Suppress certain warnings
# # warnings.filterwarnings("ignore", message="Cannot parse header or footer so it will be ignored")
# # warnings.filterwarnings("ignore", message="Data Validation extension is not supported and will be removed")
# # warnings.filterwarnings("ignore", message="Conditional Formatting extension is not supported and will be removed")
# # warnings.filterwarnings("ignore", message="Unknown extension is not supported and will be removed")
# # warnings.filterwarnings("ignore", message="Slicer List extension is not supported and will be removed")

# CHROMA_PATH = "./chroma_db"

# def main(data_path):
#     parser = argparse.ArgumentParser()
#     parser.add_argument("--reset", action="store_true", help="Reset the database.")
#     args = parser.parse_args()

#     if args.reset:
#         print("Clearing Database")
#         clear_database()

#     print("Loading documents")
#     documents = load_documents(data_path)
#     print(f"Loaded {len(documents)} documents")
#     chunks = split_documents(documents)
#     print(f"Split into {len(chunks)} chunks")
#     add_to_chroma(chunks)
#     print(f"added {len(chunks)} into Chroma DB successfully :)")

# def load_documents(data_path):
#     all_documents = []

#     print("Loading PDF documents")
#     pdf_documents = load_pdfs_recursive(data_path)
#     all_documents.extend(pdf_documents)
#     print(f"Loaded {len(pdf_documents)} PDF documents")

#     print("Loading TXT documents")
#     txt_documents = load_txt_files_recursive(data_path)
#     all_documents.extend(txt_documents)
#     print(f"Loaded {len(txt_documents)} TXT documents")

#     print("Loading DOCX documents")
#     docx_documents = load_docx_files_recursive(data_path)
#     all_documents.extend(docx_documents)
#     print(f"Loaded {len(docx_documents)} DOCX documents")

#     print("Loading PPTX documents")
#     pptx_documents = load_pptx_files_recursive(data_path)
#     all_documents.extend(pptx_documents)
#     print(f"Loaded {len(pptx_documents)} PPTX documents")

#     return all_documents

# def load_pdfs_recursive(directory):
#     pdf_documents = []
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.lower().endswith(".pdf"):
#                 file_path = os.path.join(root, file)
#                 loader = PyPDFLoader(file_path)
#                 try:
#                     docs = loader.load()
#                     pdf_documents.extend(docs)
#                 except Exception as e:
#                     print(f"Error loading PDF file {file_path}: {e}")
#     return pdf_documents

# def load_txt_files_recursive(directory):
#     txt_documents = []
#     encodings = ['iso-8859-1', 'latin-1', 'utf-8']
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.endswith(".txt"):
#                 file_path = os.path.join(root, file)
#                 content = ""
#                 for encoding in encodings:
#                     try:
#                         with open(file_path, 'r', encoding=encoding) as f:
#                             content = f.read()
#                         break
#                     except (UnicodeDecodeError, IOError) as e:
#                         print(f"Failed to read {file_path} with encoding {encoding}: {e}")
#                 txt_documents.append(
#                     Document(
#                         page_content=content,
#                         metadata={"source": file_path, "title": os.path.splitext(file)[0]}
#                     )
#                 )
#     return txt_documents

# def load_docx_files_recursive(directory):
#     docx_documents = []
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.endswith(".docx"):
#                 file_path = os.path.join(root, file)
#                 try:
#                     doc = DocxDocument(file_path)
#                     content = "\n".join([para.text for para in doc.paragraphs])
#                     docx_documents.append(
#                         Document(
#                             page_content=content,
#                             metadata={"source": file_path, "title": os.path.splitext(file)[0]}
#                         )
#                     )
#                 except Exception as e:
#                     print(f"Error loading DOCX file {file_path}: {e}")
#     return docx_documents

# def load_pptx_files_recursive(directory):
#     pptx_documents = []
#     for root, _, files in os.walk(directory):
#         for file in files:
#             if file.endswith(".pptx"):
#                 file_path = os.path.join(root, file)
#                 try:
#                     prs = Presentation(file_path)
#                     text = []
#                     for slide in prs.slides:
#                         for shape in slide.shapes:
#                             if shape.has_text_frame:
#                                 text.append(shape.text)
#                     pptx_documents.append(
#                         Document(
#                             page_content='\n'.join(text),
#                             metadata={"source": file_path, "title": os.path.splitext(file)[0]}
#                         )
#                     )
#                 except Exception as e:
#                     print(f"Error loading PPTX file {file_path}: {e}")
#     return pptx_documents

# # def split_documents(documents):
# #     # If you have a specialized text splitter, adapt or install accordingly.
# #     text_splitter = RecursiveCharacterTextSplitter(
# #         chunk_size=1000,
# #         chunk_overlap=10,
# #         length_function=len,
# #         is_separator_regex=False
# #     )
# #     chunks = text_splitter.split_documents(documents)
# #     return chunks

# def split_documents(documents):
#     """
#     Split documents by sentences using SpaCy, 
#     respecting semantic boundaries rather than raw character counts.
#     """
#     # Create the SpaCy splitter
#     text_splitter = SpacyTextSplitter(
#         pipeline='en_core_web_sm',  # the model name you downloaded
#         chunk_size=1000,            # approximate number of characters (not tokens)
#         chunk_overlap=50            # overlap in characters between chunks
#     )

#     all_chunks = []
#     for doc in documents:
#         # Split into smaller text pieces
#         text_splits = text_splitter.split_text(doc.page_content)
        
#         # Convert each piece back into a Document with the same metadata
#         for chunk_str in text_splits:
#             chunk_doc = Document(
#                 page_content=chunk_str,
#                 metadata=doc.metadata
#             )
#             all_chunks.append(chunk_doc)

#     return all_chunks

# def get_database():
#     """
#     Chroma without distance_metric parameter means it uses the default metric
#     (commonly L2 distance). 
#     """
#     return Chroma(
#         persist_directory=CHROMA_PATH,
#         embedding_function=get_embedding_function()
#     )

# def add_to_chroma(chunks):
#     chunks_with_ids = calculate_chunk_ids(chunks)
#     db = get_database()
#     db.add_documents(chunks_with_ids)

# def calculate_chunk_ids(chunks):
#     for chunk in chunks:
#         chunk.metadata["id"] = f'{chunk.metadata.get("source")}'
#     return chunks

# def clear_database():
#     if os.path.exists(CHROMA_PATH):
#         shutil.rmtree(CHROMA_PATH)
#         print("Database cleared")

# if __name__ == "__main__":
#     DATA_PATH = "./sampledata"
#     main(DATA_PATH)



#IMAGE IMPLEMENTATION:





import argparse
import os
import time
import shutil
import warnings
from langchain_community.document_loaders import PyPDFDirectoryLoader
from langchain_community.document_loaders import PyPDFLoader
from langchain.schema.document import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from get_embedding import get_embedding_function
from langchain.text_splitter import SpacyTextSplitter

from docx import Document as DocxDocument
from langchain_chroma import Chroma
from pptx import Presentation
import pandas as pd

# Optional: Suppress certain warnings
# warnings.filterwarnings("ignore", message="Cannot parse header or footer so it will be ignored")
# warnings.filterwarnings("ignore", message="Data Validation extension is not supported and will be removed")
# warnings.filterwarnings("ignore", message="Conditional Formatting extension is not supported and will be removed")
# warnings.filterwarnings("ignore", message="Unknown extension is not supported and will be removed")
# warnings.filterwarnings("ignore", message="Slicer List extension is not supported and will be removed")

CHROMA_PATH = "./chroma_db"

def main(data_path):
    parser = argparse.ArgumentParser()
    parser.add_argument("--reset", action="store_true", help="Reset the database.")
    args = parser.parse_args()

    if args.reset:
        print("Clearing Database")
        clear_database()

    print("Loading documents")
    documents = load_documents(data_path)
    print(f"Loaded {len(documents)} documents")
    chunks = split_documents(documents)
    print(f"Split into {len(chunks)} chunks")
    add_to_chroma(chunks)
    print(f"added {len(chunks)} into Chroma DB successfully :)")

def load_documents(data_path):
    all_documents = []

    print("Loading IMAGE documents")
    image_docs = load_image_files_recursive(data_path)
    all_documents.extend(image_docs)
    print(f"Loaded {len(image_docs)} image documents")

    print("Loading PDF documents")
    pdf_documents = load_pdfs_recursive(data_path)
    all_documents.extend(pdf_documents)
    print(f"Loaded {len(pdf_documents)} PDF documents")

    print("Loading TXT documents")
    txt_documents = load_txt_files_recursive(data_path)
    all_documents.extend(txt_documents)
    print(f"Loaded {len(txt_documents)} TXT documents")

    print("Loading DOCX documents")
    docx_documents = load_docx_files_recursive(data_path)
    all_documents.extend(docx_documents)
    print(f"Loaded {len(docx_documents)} DOCX documents")

    print("Loading PPTX documents")
    pptx_documents = load_pptx_files_recursive(data_path)
    all_documents.extend(pptx_documents)
    print(f"Loaded {len(pptx_documents)} PPTX documents")

    return all_documents

def load_pdfs_recursive(directory):
    pdf_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".pdf"):
                file_path = os.path.join(root, file)
                loader = PyPDFLoader(file_path)
                try:
                    docs = loader.load()
                    pdf_documents.extend(docs)
                except Exception as e:
                    print(f"Error loading PDF file {file_path}: {e}")
    return pdf_documents



def load_image_files_recursive(directory):
    """
    Recursively load all image files from a directory and return them as Documents.
    Each Document will have minimal text, but the main content is the image path in metadata.
    """
    exts = [".png", ".jpg", ".jpeg", ".bmp", ".gif", ".tiff"]  # extend as needed
    image_docs = []
    for root, _, files in os.walk(directory):
        for file in files:
            if any(file.lower().endswith(ext) for ext in exts):
                file_path = os.path.join(root, file)
                # For images, we might store a placeholder 'page_content'
                # Because the main retrieval is from the CLIP image embedding
                # But you can store a caption if you have it
                doc = Document(
                    page_content=f"Image file: {file_path}",
                    metadata={"source": file_path, "title": os.path.splitext(file)[0]}
                )
                image_docs.append(doc)
    return image_docs

def load_txt_files_recursive(directory):
    txt_documents = []
    encodings = ['iso-8859-1', 'latin-1', 'utf-8']
    for root, _, files in os.walk(directory):
        for file in files:
            if file.endswith(".txt"):
                file_path = os.path.join(root, file)
                content = ""
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read()
                        break
                    except (UnicodeDecodeError, IOError) as e:
                        print(f"Failed to read {file_path} with encoding {encoding}: {e}")
                txt_documents.append(
                    Document(
                        page_content=content,
                        metadata={"source": file_path, "title": os.path.splitext(file)[0]}
                    )
                )
    return txt_documents

def load_docx_files_recursive(directory):
    docx_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.endswith(".docx"):
                file_path = os.path.join(root, file)
                try:
                    doc = DocxDocument(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                    docx_documents.append(
                        Document(
                            page_content=content,
                            metadata={"source": file_path, "title": os.path.splitext(file)[0]}
                        )
                    )
                except Exception as e:
                    print(f"Error loading DOCX file {file_path}: {e}")
    return docx_documents

def load_pptx_files_recursive(directory):
    pptx_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.endswith(".pptx"):
                file_path = os.path.join(root, file)
                try:
                    prs = Presentation(file_path)
                    text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if shape.has_text_frame:
                                text.append(shape.text)
                    pptx_documents.append(
                        Document(
                            page_content='\n'.join(text),
                            metadata={"source": file_path, "title": os.path.splitext(file)[0]}
                        )
                    )
                except Exception as e:
                    print(f"Error loading PPTX file {file_path}: {e}")
    return pptx_documents

# def split_documents(documents):
#     # If you have a specialized text splitter, adapt or install accordingly.
#     text_splitter = RecursiveCharacterTextSplitter(
#         chunk_size=1000,
#         chunk_overlap=10,
#         length_function=len,
#         is_separator_regex=False
#     )
#     chunks = text_splitter.split_documents(documents)
#     return chunks

def split_documents(documents):
    """
    Split documents by sentences using SpaCy, 
    respecting semantic boundaries rather than raw character counts.
    """
    # Create the SpaCy splitter
    text_splitter = SpacyTextSplitter(
        pipeline='en_core_web_sm',  # the model name you downloaded
        chunk_size=1000,            # approximate number of characters (not tokens)
        chunk_overlap=50            # overlap in characters between chunks
    )

    all_chunks = []
    for doc in documents:
        # Split into smaller text pieces
        text_splits = text_splitter.split_text(doc.page_content)
        
        # Convert each piece back into a Document with the same metadata
        for chunk_str in text_splits:
            chunk_doc = Document(
                page_content=chunk_str,
                metadata=doc.metadata
            )
            all_chunks.append(chunk_doc)

    return all_chunks

def get_database():
    """
    Chroma without distance_metric parameter means it uses the default metric
    (commonly L2 distance). 
    """
    return Chroma(
        persist_directory=CHROMA_PATH,
        embedding_function=get_embedding_function()
    )

def add_to_chroma(chunks):
    chunks_with_ids = calculate_chunk_ids(chunks)
    db = get_database()
    db.add_documents(chunks_with_ids)

def calculate_chunk_ids(chunks):
    for chunk in chunks:
        chunk.metadata["id"] = f'{chunk.metadata.get("source")}'
    return chunks

def clear_database():
    if os.path.exists(CHROMA_PATH):
        shutil.rmtree(CHROMA_PATH)
        print("Database cleared")

if __name__ == "__main__":
    DATA_PATH = "./sampledata"
    main(DATA_PATH)
