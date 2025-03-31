
import argparse
import os
import time
import shutil
import warnings
from langchain_community.document_loaders import PyPDFLoader
from langchain.schema import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from get_embedding import get_embedding_function
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
from langchain_chroma import Chroma

#  suppress warnings 
warnings.filterwarnings("ignore")

CHROMA_PATH = "./chroma"  # not used in per-patient function
DATA_PATH = "./sampledata"

# -----------------------
# loader functions for each file type:
# -----------------------

def load_pdfs_recursive(directory):
    pdf_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".pdf"):
                file_path = os.path.join(root, file)
                loader = PyPDFLoader(file_path)
                try:
                    docs = loader.load()  # returns list of Document objects
                    pdf_documents.extend(docs)
                except Exception as e:
                    print(f"Error loading PDF file {file_path}: {e}")
    return pdf_documents

def load_txt_files_recursive(directory):
    txt_documents = []
    encodings = ['utf-8', 'iso-8859-1', 'latin-1']
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".txt"):
                file_path = os.path.join(root, file)
                content = ""
                for encoding in encodings:
                    try:
                        with open(file_path, 'r', encoding=encoding) as f:
                            content = f.read()
                        break
                    except Exception as e:
                        print(f"Error reading TXT file {file_path} with encoding {encoding}: {e}")
                txt_documents.append(Document(page_content=content, metadata={
                    "source": file_path,
                    "title": os.path.splitext(file)[0]
                }))
                print(txt_documents)
    return txt_documents

def load_docx_files_recursive(directory):
    docx_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".docx"):
                file_path = os.path.join(root, file)
                try:
                    doc = DocxDocument(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                    docx_documents.append(Document(page_content=content, metadata={
                        "source": file_path,
                        "title": os.path.splitext(file)[0]
                    }))
                except Exception as e:
                    print(f"Error loading DOCX file {file_path}: {e}")
    return docx_documents

def load_pptx_files_recursive(directory):
    pptx_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith(".pptx"):
                file_path = os.path.join(root, file)
                try:
                    prs = Presentation(file_path)
                    content = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                content.append(shape.text)
                    pptx_documents.append(Document(page_content="\n".join(content), metadata={
                        "source": file_path,
                        "title": os.path.splitext(file)[0]
                    }))
                except Exception as e:
                    print(f"Error loading PPTX file {file_path}: {e}")
    return pptx_documents

def load_excel_files_recursive(directory):
    excel_documents = []
    for root, _, files in os.walk(directory):
        for file in files:
            if file.lower().endswith((".xlsx", ".xls")):
                file_path = os.path.join(root, file)
                try:
                    df = pd.read_excel(file_path)
                    content = df.to_string(index=False)
                    excel_documents.append(Document(page_content=content, metadata={
                        "source": file_path,
                        "title": os.path.splitext(file)[0]
                    }))
                except Exception as e:
                    print(f"Error loading Excel file {file_path}: {e}")
    return excel_documents
def get_database():
    return Chroma(
        persist_directory=CHROMA_PATH,
        embedding_function=get_embedding_function()  # MPNet-based
    )
def load_documents(data_path):
    """Load documents from all supported file types in the directory."""
    all_documents = []
    print("Loading PDF documents")
    pdf_docs = load_pdfs_recursive(data_path)
    all_documents.extend(pdf_docs)
    print(f"Loaded {len(pdf_docs)} PDF documents")

    print("Loading TXT documents")
    txt_docs = load_txt_files_recursive(data_path)
    all_documents.extend(txt_docs)
    print(f"Loaded {len(txt_docs)} TXT documents")

    print("Loading DOCX documents")
    docx_docs = load_docx_files_recursive(data_path)
    all_documents.extend(docx_docs)
    print(f"Loaded {len(docx_docs)} DOCX documents")

    print("Loading PPTX documents")
    pptx_docs = load_pptx_files_recursive(data_path)
    all_documents.extend(pptx_docs)
    print(f"Loaded {len(pptx_docs)} PPTX documents")

    print("Loading Excel documents")
    excel_docs = load_excel_files_recursive(data_path)
    all_documents.extend(excel_docs)
    print(f"Loaded {len(excel_docs)} Excel documents")

    return all_documents

# -----------------------
#  Contextual Chunking Function (Word-based with overlap)
# -----------------------

def contextual_split_documents(documents, max_words=150, overlap_words=20):
    """Split documents into chunks based on word counts, preserving full words and with overlap."""
    new_chunks = []
    for doc in documents:
        words = doc.page_content.split()
        if len(words) <= max_words:
            new_chunks.append(doc)
            continue
        start = 0
        while start < len(words):
            end = start + max_words
            chunk_words = words[start:end]
            chunk_text = " ".join(chunk_words)
            new_doc = Document(page_content=chunk_text, metadata=doc.metadata.copy())
            new_chunks.append(new_doc)
            if end >= len(words):
                break
            start = end - overlap_words  
    return new_chunks

# -----------------------
# New Function: populate_patient_vector_db
# -----------------------

def populate_patient_vector_db(patient_folder):
    """
    Populate the vector database for a single patient.
    Expects patient_folder to be the path to the patient's directory,
    e.g. "sampledata/<patientId>".
    This function will:
      - Load documents from patient_folder.
      - Use contextual chunking.
      - Create (or update) a Chroma vector store in patient_folder/chromaDB.
    """
    print(f"Populating vector DB for patient folder: {patient_folder}")
    
    #  where the vector DB will be stored:
    vector_db_folder = os.path.join(patient_folder, "chromaDB")
    if not os.path.exists(vector_db_folder):
        os.makedirs(vector_db_folder, exist_ok=True)
        print(f"Created vector DB folder: {vector_db_folder}")
    
    # load all documents from the patient folder.
    documents = load_documents(patient_folder)
    print(f"Loaded {len(documents)} documents from {patient_folder}")
    
    # no documents found, you might choose to initialize an empty DB:
    if not documents:
        print("No documents found; initializing an empty vector store.")
        # init empty Chroma store:
        vector_db = Chroma(
            persist_directory=vector_db_folder,
            embedding_function=get_embedding_function(),
        )

        return
    
    # split documents using contextual chunking.
    chunks = contextual_split_documents(documents, max_words=150, overlap_words=20)
    print(f"Split into {len(chunks)} chunks")
    
    # init the Chroma vector store for this patient.
    vector_db = Chroma(
        persist_directory=vector_db_folder,
        embedding_function=get_embedding_function(),
    )
    
  
    for chunk in chunks:
        chunk.metadata["id"] = str(time.time())
    
    vector_db.add_documents(chunks)
    print("Finished adding chunks to patient vector DB.")

# -----------------------
# Add to Chroma using global approach (if needed)
# -----------------------

def add_to_chroma(chunks):
    print("Adding chunks to global Chroma DB")
    vector_db = Chroma(
        persist_directory=CHROMA_PATH,
        embedding_function=get_embedding_function(),
    )
    for chunk in chunks:
        chunk.metadata["id"] = str(time.time())
    vector_db.add_documents(chunks)
    print("Finished adding chunks to global Chroma DB")

def clear_database():
    if os.path.exists(CHROMA_PATH):
        print("Clearing Database")
        shutil.rmtree(CHROMA_PATH)
        print("Database cleared")
    else:
        print("Database directory does not exist")

def main(data_path):
    parser = argparse.ArgumentParser()
    parser.add_argument("--reset", action="store_true", help="Reset the database.")
    args = parser.parse_args()
    if args.reset:
        clear_database()

    print("Loading documents from", data_path)
    documents = load_documents(data_path)
    print(f"Loaded {len(documents)} documents")
    
    chunks = contextual_split_documents(documents, max_words=150, overlap_words=20)
    print(f"Split into {len(chunks)} chunks")
    
    add_to_chroma(chunks)
    print("Done populating global vector database.")

if __name__ == "__main__":
    main(DATA_PATH)
